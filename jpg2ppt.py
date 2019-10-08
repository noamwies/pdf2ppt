import argparse
import glob

from pptx import Presentation
from pptx.util import Inches
from tqdm import tqdm


BLANK = 6


def create_presentation_from_images(name, images):
	prs = Presentation()
	left = top = Inches(0)
	slide_layout = prs.slide_layouts[BLANK]
	for img_path in tqdm(images):
		slide = prs.slides.add_slide(slide_layout)
		pic = slide.shapes.add_picture(img_path, left, top)
	prs.save(name)


def get_config():
    parser = argparse.ArgumentParser(description='Pdf generator.')
    parser.add_argument('--output_path', type=str, required=True)
    parser.add_argument('--images_regex', type=str, default='*.jpg')
    return parser.parse_args()


def main():
	config = get_config()
	images = sorted(glob.glob(config.images_regex))
	create_presentation_from_images(config.output_path, images)

	
if __name__ == '__main__':
	main()
